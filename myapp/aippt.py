from pptx import Presentation

# Create a new PowerPoint presentation
presentation = Presentation()

# Slide 1 - Title slide
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]
title_1.text = "The History of Artificial Intelligence"
subtitle_1.text = "An overview of the key milestones"

# Slide 2 - Introduction
slide_2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_2 = slide_2.shapes.title
content_2 = slide_2.placeholders[1]
title_2.text = "Introduction"
content_2.text = "Artificial Intelligence (AI) refers to the development of computer systems capable of performing tasks that typically require human intelligence. The history of AI dates back to the mid-20th century and has witnessed remarkable advancements."

# Slide 3 - Early Developments
slide_3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_3 = slide_3.shapes.title
content_3 = slide_3.placeholders[1]
title_3.text = "Early Developments"
content_3.text = "In the 1950s, AI pioneers like Alan Turing and John McCarthy laid the foundations for AI research. Turing proposed the concept of a 'universal machine' capable of simulating any other machine, while McCarthy organized the Dartmouth Conference, widely considered the birth of AI."

# Slide 4 - AI Winter and Resurgence
slide_4 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_4 = slide_4.shapes.title
content_4 = slide_4.placeholders[1]
title_4.text = "AI Winter and Resurgence"
content_4.text = "The late 20th century saw an 'AI winter' with reduced funding and interest due to the inability to fulfill early expectations. However, breakthroughs in the 21st century, such as deep learning and big data, fueled the resurgence of AI and its integration into various domains."

# Slide 5 - Current Applications and Future Prospects
slide_5 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_5 = slide_5.shapes.title
content_5 = slide_5.placeholders[1]
title_5.text = "Current Applications and Future Prospects"
content_5.text = "AI has transformed industries, enabling advancements in autonomous vehicles, natural language processing, computer vision, and more. With ongoing research in areas like explainable AI and ethical considerations, the future of AI holds immense potential."

# Save the presentation-
presentation.save("history_of_ai.pptx")


